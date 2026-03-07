"""
Multi-format document parser for Saarthi.AI
Supports: PDF, DOCX, DOC, TXT, Images, HTML, XLSX, XLS, PPTX, CSV, RTF, ODT, MD

Production pipeline uses Amazon Textract for PDF and image OCR.
"""

import os
import sys
import logging
import boto3
from typing import Optional

logger = logging.getLogger(__name__)

# ── AWS Textract ────────────────────────────────────────────────
textract_client = None
try:
    textract_client = boto3.client("textract", region_name=os.getenv("AWS_REGION", "ap-south-1"))
except Exception as e:
    logger.warning(f"Failed to initialize Textract client: {e}")

# ── DOCX ───────────────────────────────────────────────────────
try:
    import docx  # python-docx
except ImportError:
    docx = None

# ── HTML ───────────────────────────────────────────────────────
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

# ── Excel (XLSX, XLS) ─────────────────────────────────────────
try:
    import openpyxl  # For .xlsx files
except ImportError:
    openpyxl = None

try:
    import xlrd  # For .xls files (older Excel format)
except ImportError:
    xlrd = None

# ── PowerPoint (PPTX) ─────────────────────────────────────────
try:
    from pptx import Presentation  # python-pptx
except ImportError:
    Presentation = None

# ── CSV ───────────────────────────────────────────────────────
import csv  # Built-in, always available

# ── RTF ───────────────────────────────────────────────────────
try:
    import striprtf  # For RTF files
except ImportError:
    striprtf = None

# ── ODT (OpenDocument) ────────────────────────────────────────
try:
    from odf.opendocument import load as odt_load  # odfpy
    from odf.text import P as OdfP
except ImportError:
    try:
        # Alternative import path
        from odf import opendocument
        odt_load = opendocument.load
        from odf.text import P as OdfP
    except ImportError:
        odt_load = None
        OdfP = None


# ───────────────────────────────────────────────────────────────
# Amazon Textract extraction (Production)
# ───────────────────────────────────────────────────────────────

def _extract_lines_from_blocks(blocks) -> str:
    """Helper to concatenate LINE blocks into a single string."""
    text_lines = []
    for block in blocks or []:
        if block.get("BlockType") == "LINE":
            text = (block.get("Text") or "").strip()
            if text:
                text_lines.append(text)
    return "\n".join(text_lines)


def extract_text_from_image(image_bytes: bytes) -> str:
    """
    Extract text from an image (bytes) using Textract DetectDocumentText.
    Used for PNG/JPG/JPEG/TIFF/BMP document images.
    """
    if not textract_client:
        raise Exception("Textract client not initialized. Check AWS credentials and region.")

    logger.info("Starting Textract DetectDocumentText for image bytes (size=%d)", len(image_bytes))

    try:
        response = textract_client.detect_document_text(
            Document={"Bytes": image_bytes}
        )
        blocks = response.get("Blocks", []) or []
        logger.info("Textract image detection returned %d blocks", len(blocks))

        extracted_text = _extract_lines_from_blocks(blocks)
        logger.info("Textract image extraction completed: %d characters", len(extracted_text))

        if not extracted_text:
            raise Exception("Document text extraction failed: No text found in image")

        return extracted_text

    except Exception as e:
        logger.error("Textract image extraction error: %s", str(e), exc_info=True)
        raise


def extract_text_from_pdf(bucket_name: str, object_key: str) -> str:
    """
    Extract text from a PDF stored in S3 using the asynchronous Textract
    StartDocumentTextDetection / GetDocumentTextDetection pipeline.
    """
    if not textract_client:
        raise Exception("Textract client not initialized. Check AWS credentials and region.")

    logger.info(
        "Starting Textract async PDF job for: s3://%s/%s",
        bucket_name,
        object_key,
    )

    try:
        # Start async Textract job for the PDF in S3
        job = textract_client.start_document_text_detection(
            DocumentLocation={
                "S3Object": {
                    "Bucket": bucket_name,
                    "Name": object_key,
                }
            }
        )

        job_id = job.get("JobId")
        if not job_id:
            raise Exception("Textract did not return a JobId for PDF document.")

        logger.info("Textract async PDF job started. JobId=%s", job_id)

        # Poll job status with a 180 second max wait
        import time

        max_wait_seconds = 180
        poll_interval_seconds = 3
        start_time = time.time()

        response = None
        while True:
            response = textract_client.get_document_text_detection(JobId=job_id)
            status = response.get("JobStatus")
            logger.info("Waiting for Textract job %s status: %s", job_id, status)

            if status == "SUCCEEDED":
                break
            if status == "FAILED":
                message = response.get("StatusMessage") or "Textract PDF job failed."
                raise Exception(f"Document text extraction failed: {message}")

            if time.time() - start_time > max_wait_seconds:
                raise Exception("Textract PDF job timed out after 180 seconds.")

            time.sleep(poll_interval_seconds)

        # Collect all blocks (handle pagination via NextToken)
        blocks = list((response or {}).get("Blocks", []) or [])
        next_token = (response or {}).get("NextToken")
        while next_token:
            page = textract_client.get_document_text_detection(JobId=job_id, NextToken=next_token)
            blocks.extend(page.get("Blocks", []) or [])
            next_token = page.get("NextToken")

        logger.info("Textract async PDF extraction returned %d blocks", len(blocks))

        extracted_text = _extract_lines_from_blocks(blocks)
        logger.info("Textract extraction completed: %d characters", len(extracted_text))

        if not extracted_text:
            raise Exception("Document text extraction failed: No text found in PDF")

        return extracted_text

    except Exception as e:
        logger.error("Textract PDF extraction error: %s", str(e), exc_info=True)
        raise


def extract_text_with_textract(bucket_name: str, object_key: str) -> str:
    """
    High-level Textract router for S3 documents.

    - PDFs -> asynchronous StartDocumentTextDetection / GetDocumentTextDetection
    - Images (PNG/JPG/JPEG/TIFF/BMP) -> DetectDocumentText on image bytes
    """
    if not textract_client:
        raise Exception("Textract client not initialized. Check AWS credentials and region.")

    normalized_key = (object_key or "").strip()
    basename = os.path.basename(normalized_key).lower()
    extension = basename.rsplit(".", 1)[-1] if "." in basename else ""

    image_exts = {"png", "jpg", "jpeg", "tiff", "tif", "bmp"}

    logger.info(
        "Processing document with Textract router: s3://%s/%s filename=%s extension=%s",
        bucket_name,
        object_key,
        basename or "unknown",
        extension or "unknown",
    )

    # Validation: only allow PDF + image formats here
    if extension == "pdf":
        logger.info("Detected extension: pdf; using Textract async PDF pipeline")
        return extract_text_from_pdf(bucket_name, object_key)
    elif extension in image_exts:
        logger.info("Detected extension: %s; using Textract image pipeline", extension)
        # Download image bytes from S3 then run DetectDocumentText on bytes
        s3_client = boto3.client("s3")
        obj = s3_client.get_object(Bucket=bucket_name, Key=object_key)
        image_bytes = obj["Body"].read()
        text = extract_text_from_image(image_bytes)
        logger.info("Textract image pipeline extracted %d characters", len(text))
        return text
    else:
        logger.warning("Unsupported document type for Textract router: extension=%s", extension)
        raise Exception("Unsupported document type. Supported formats: PDF, PNG, JPG, JPEG, TIFF, BMP")


# ───────────────────────────────────────────────────────────────
# File type detection
# ───────────────────────────────────────────────────────────────

def detect_file_type(filename: str) -> str:
    """Detect file type from extension."""
    if not filename:
        return "unknown"
    ext = filename.lower().rsplit(".", 1)[-1]
    
    # Documents
    if ext == "pdf":
        return "pdf"
    elif ext in ("docx", "doc"):
        return "docx"
    elif ext in ("xlsx", "xls"):
        return "excel"
    elif ext in ("pptx", "ppt"):
        return "powerpoint"
    elif ext == "odt":
        return "odt"
    elif ext == "rtf":
        return "rtf"
    # Text formats
    elif ext == "txt":
        return "txt"
    elif ext in ("md", "markdown"):
        return "markdown"
    elif ext == "csv":
        return "csv"
    # Web formats
    elif ext in ("html", "htm"):
        return "html"
    # Images
    elif ext in ("png", "jpg", "jpeg", "bmp", "tiff", "tif", "webp", "gif", "svg"):
        return "image"
    else:
        return "unknown"


# ───────────────────────────────────────────────────────────────
# Main extraction
# ───────────────────────────────────────────────────────────────

def extract_text(file_path: str, filename: str = None) -> str:
    """
    Extract text from a file.  Automatically detects the format
    from *filename* and applies the appropriate parser.

    For images and scanned PDFs, applies image preprocessing
    before OCR.

    Args:
        file_path: Path to the file on disk.
        filename:  Original filename (used for type detection).

    Returns:
        Extracted text string (may be empty on failure).
    """
    if not filename:
        filename = os.path.basename(file_path)

    file_type = detect_file_type(filename)

    try:
        # ── PDF ────────────────────────────────────────────────
        if file_type == "pdf":
            # PDFs in production are handled via Textract on S3 objects.
            # Local file-path based PDF parsing is intentionally not implemented
            # to avoid binary dependencies in Lambda.
            raise ValueError(
                "Local PDF extraction is not supported here. Use the S3 + Textract "
                "pipeline (extract_text_with_textract) for PDFs."
            )

        # ── DOCX ───────────────────────────────────────────────
        elif file_type == "docx":
            if not docx:
                raise ImportError("python-docx is not installed.")
            doc = docx.Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs)

        # ── TXT ────────────────────────────────────────────────
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        # ── IMAGE ──────────────────────────────────────────────
        elif file_type == "image":
            # Image OCR in production is handled via Textract in extract_text_with_textract
            # using S3 objects. This local helper does not implement image OCR anymore.
            print(
                "Image OCR via local file path is not supported here. "
                "Use the S3 + Textract pipeline (extract_text_with_textract).",
                file=sys.stderr,
            )
            return ""

        # ── HTML ───────────────────────────────────────────────
        elif file_type == "html":
            if not BeautifulSoup:
                raise ImportError("beautifulsoup4 is not installed.")
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                return soup.get_text(separator="\n", strip=True)

        # ── EXCEL (XLSX, XLS) ──────────────────────────────────
        elif file_type == "excel":
            ext = filename.lower().rsplit(".", 1)[-1]
            text_lines = []
            
            if ext == "xlsx":
                if not openpyxl:
                    raise ImportError("openpyxl is not installed. Install with: pip install openpyxl")
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_lines.append(f"\n--- Sheet: {sheet_name} ---\n")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                        if row_text.strip():
                            text_lines.append(row_text)
            elif ext == "xls":
                if not xlrd:
                    raise ImportError("xlrd is not installed. Install with: pip install xlrd")
                workbook = xlrd.open_workbook(file_path)
                for sheet in workbook.sheets():
                    text_lines.append(f"\n--- Sheet: {sheet.name} ---\n")
                    for row_idx in range(sheet.nrows):
                        row = sheet.row(row_idx)
                        row_text = "\t".join(str(cell.value) for cell in row)
                        if row_text.strip():
                            text_lines.append(row_text)
            
            return "\n".join(text_lines)

        # ── POWERPOINT (PPTX) ───────────────────────────────────
        elif file_type == "powerpoint":
            if not Presentation:
                raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
            prs = Presentation(file_path)
            text_lines = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_lines.append(f"\n--- Slide {slide_num} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines.append(shape.text)
            return "\n".join(text_lines)

        # ── CSV ────────────────────────────────────────────────
        elif file_type == "csv":
            text_lines = []
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text_lines.append("\t".join(row))
            except UnicodeDecodeError:
                # Try with different encoding
                with open(file_path, "r", encoding="latin-1") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text_lines.append("\t".join(row))
            return "\n".join(text_lines)

        # ── RTF ────────────────────────────────────────────────
        elif file_type == "rtf":
            if not striprtf:
                raise ImportError("striprtf is not installed. Install with: pip install striprtf")
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf_text = f.read()
                from striprtf.striprtf import rtf_to_text
                return rtf_to_text(rtf_text)

        # ── ODT (OpenDocument) ──────────────────────────────────
        elif file_type == "odt":
            if not odt_load or not OdfP:
                raise ImportError("odfpy is not installed. Install with: pip install odfpy")
            doc = odt_load(file_path)
            text_lines = []
            for paragraph in doc.getElementsByType(OdfP):
                text = paragraph.__str__()
                if text.strip():
                    text_lines.append(text)
            return "\n".join(text_lines)

        # ── MARKDOWN ───────────────────────────────────────────
        elif file_type == "markdown":
            # Markdown is plain text, just read it
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        else:
            # For unknown types, try to read as text (might work for some files)
            print(f"Warning: Unknown file type '{file_type}' for {filename}, attempting to read as text", file=sys.stderr)
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            except Exception:
                raise ValueError(f"Unsupported file type for {filename}. Supported: PDF, DOCX, DOC, TXT, XLSX, XLS, PPTX, CSV, RTF, ODT, MD, HTML, Images")

    except Exception as e:
        print(f"Error extracting text from {file_path} (Type: {file_type}): {e}", file=sys.stderr)
        return ""
